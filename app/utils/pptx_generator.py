from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import streamlit as st


def set_background_color(slide, rgb_color):
    """
    Set the background color of a slide.

    Args:
        slide (pptx.slide.Slide): Slide object to modify.
        rgb_color (tuple): RGB color as a tuple of three integers (R, G, B).
    """
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)


def generate_slides(
    slide_content,
    slide_num=None,
    filename="presentation.pptx",
    layout=1,
    background_color=(255, 255, 255),
    font_color=(0,0,0)
):
    """
    Generate a PowerPoint presentation file from given slide content and
    provide a Streamlit download button.

    Args:
        slide_content (list of dict): List of slides, each dict contains 'title' and 'content' (list of strings).
        slide_num (int, optional): Number of slides to generate. Defaults to length of slide_content.
        filename (str): Name of the output pptx file.
        layout (int): PPTX slide layout index to use.
        background_color (tuple): RGB color tuple for slide background.
    """
    prs = Presentation()

    if slide_num is None:
        slide_num = len(slide_content)

    for i in range(slide_num):
        slide = prs.slides.add_slide(prs.slide_layouts[layout])
        set_background_color(slide, background_color)

        # Set slide title
        title_shape = slide.shapes.title
        title_shape.text = slide_content[i].get("title", "")
        # Set title font color to black and font size (optional)
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.color.rgb = RGBColor(*font_color)
        title_paragraph.font.size = Pt(36)  # example size, adjust as needed

        # Set slide body content
        body_shape = slide.shapes.placeholders[1]
        text_frame = body_shape.text_frame

        content = slide_content[i].get("content", [])
        if content:
            # Set first bullet point text and style
            text_frame.text = content[0]
            p0 = text_frame.paragraphs[0]
            p0.font.size = Pt(20)
            p0.font.color.rgb = RGBColor(*font_color)  # black font color

            # Add remaining bullet points
            for bullet_text in content[1:]:
                p = text_frame.add_paragraph()
                p.text = bullet_text
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(*font_color)  # black font color
                p.level = 0  # same indentation level
        else:
            text_frame.text = ""

    # Save the presentation file
    file_path = f"data\{filename}"
    prs.save(file_path)

    # # Provide download button in Streamlit
    # with open(file_path, "rb") as file:
    #     pptx_data = file.read()

    # st.download_button(
    #     label="Download PPTX",
    #     data=pptx_data,
    #     file_name=filename,
    #     mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    # )


if __name__ == "__main__":
    sample_slides = [
        {
            "title": "History of Indonesia",
            "content": [
                "Ancient civilizations and early trade routes in Southeast Asia",
                "Colonial history: Dutch colonization and its impact",
                "Struggle for independence: Key figures and events",
                "Post-independence developments and political changes",
                "Current political structure and democracy",
            ],
        }
    ]

    generate_slides(slide_content=sample_slides)
